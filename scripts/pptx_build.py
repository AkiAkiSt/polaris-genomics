"""Build the POLARIS 10-minute deck: dark sandwich, figure-driven, minimal text, transitions."""
from __future__ import annotations
import sys, json, pathlib, copy
sys.path.insert(0, str(pathlib.Path(__file__).resolve().parents[1]))
from pptx import Presentation
from pptx.util import Inches as I, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.enum.shapes import MSO_SHAPE
RR = MSO_SHAPE.ROUNDED_RECTANGLE
from pptx.oxml.ns import qn
from PIL import Image
from polaris import config

FIG = config.FIG; TAB = config.TAB
NAVY = RGBColor(0x0b, 0x1f, 0x3a); INK = RGBColor(0x10, 0x26, 0x3f)
TEAL = RGBColor(0x1f, 0x9e, 0x89); AMBER = RGBColor(0xf0, 0xa2, 0x02)
SLATE = RGBColor(0x5b, 0x6b, 0x7b); MIST = RGBColor(0xdc, 0xe6, 0xef)
WHITE = RGBColor(0xff, 0xff, 0xff); CRIM = RGBColor(0xb3, 0x12, 0x3a)
m = json.load(open(TAB / "polaris_metrics.json")); cad = json.load(open(TAB / "cad_metrics.json"))
rob = json.load(open(TAB / "robustness.json"))

prs = Presentation(); prs.slide_width = I(13.333); prs.slide_height = I(7.5)
BLANK = prs.slide_layouts[6]
SW, SH = prs.slide_width, prs.slide_height


def bg(slide, color):
    slide.background.fill.solid(); slide.background.fill.fore_color.rgb = color


def textbox(slide, x, y, w, h, lines, align=PP_ALIGN.LEFT, anchor=MSO_ANCHOR.TOP):
    tb = slide.shapes.add_textbox(I(x), I(y), I(w), I(h)); tf = tb.text_frame
    tf.word_wrap = True; tf.vertical_anchor = anchor
    tf.margin_left = tf.margin_right = tf.margin_top = tf.margin_bottom = 0
    for i, (txt, sz, col, bold, font) in enumerate(lines):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph(); p.alignment = align
        r = p.add_run(); r.text = txt; f = r.font
        f.size = Pt(sz); f.bold = bold; f.color.rgb = col; f.name = font
    return tb


def fit_image(slide, path, box_x, box_y, box_w, box_h):
    w, h = Image.open(path).size; ar = w / h; bar = box_w / box_h
    if ar > bar:
        iw = box_w; ih = box_w / ar
    else:
        ih = box_h; iw = box_h * ar
    x = box_x + (box_w - iw) / 2; y = box_y + (box_h - ih) / 2
    slide.shapes.add_picture(str(path), I(x), I(y), I(iw), I(ih))


def transition(slide, kind="fade"):
    xml = ('<p:transition xmlns:p="http://schemas.openxmlformats.org/presentationml/2006/main" '
           'xmlns:p14="http://schemas.microsoft.com/office/powerpoint/2010/main" spd="med" p14:dur="600">'
           f'<p:{kind}/></p:transition>')
    from pptx.oxml import parse_xml
    el = parse_xml(xml)
    sld = slide._element
    csld = sld.find(qn('p:cSld'))
    clr = sld.find(qn('p:clrMapOvr'))
    (clr if clr is not None else csld).addnext(el)


def fig_slide(title, fig, takeaway, tag=None):
    s = prs.slides.add_slide(BLANK); bg(s, WHITE)
    textbox(s, 0.55, 0.34, 12.2, 0.8, [(title, 30, NAVY, True, "Cambria")])
    if tag:
        textbox(s, 0.55, 0.95, 12.2, 0.4, [(tag, 14, TEAL, False, "Calibri")])
    fit_image(s, FIG / f"{fig}.png", 0.5, 1.35, 12.33, 5.05)
    # takeaway: teal rounded callout (no stripe), bottom
    cap = s.shapes.add_shape(RR, I(0.9), I(6.55), I(11.5), I(0.62))
    cap.fill.solid(); cap.fill.fore_color.rgb = MIST; cap.line.fill.background()
    cap.adjustments[0] = 0.5
    tf = cap.text_frame; tf.word_wrap = True; tf.margin_top = Pt(4); tf.margin_bottom = Pt(4)
    p = tf.paragraphs[0]; p.alignment = PP_ALIGN.CENTER
    r = p.add_run(); r.text = takeaway; r.font.size = Pt(15); r.font.bold = True
    r.font.color.rgb = INK; r.font.name = "Calibri"
    transition(s)
    return s

# 1 — title
s = prs.slides.add_slide(BLANK); bg(s, NAVY)
textbox(s, 1.0, 2.25, 11.3, 1.5, [("POLARIS", 72, WHITE, True, "Cambria")], PP_ALIGN.CENTER)
textbox(s, 1.4, 3.75, 10.5, 1.2,
        [("Calibrated triangulation of non-coding variant evidence", 24, MIST, False, "Calibri")], PP_ALIGN.CENTER)
textbox(s, 1.4, 4.4, 10.5, 1.0,
        [("into mechanism and target gene", 24, MIST, False, "Calibri")], PP_ALIGN.CENTER)
textbox(s, 1.4, 5.5, 10.5, 0.6,
        [("type 2 diabetes + coronary artery disease  ·  transparent  ·  reproducible from public APIs",
          15, TEAL, False, "Calibri")], PP_ALIGN.CENTER)
transition(s)

# 2 — problem + guardrails
s = prs.slides.add_slide(BLANK); bg(s, WHITE)
textbox(s, 0.55, 0.4, 12.2, 0.8, [("90% of disease variants are non-coding — the bottleneck is integration", 28, NAVY, True, "Cambria")])
textbox(s, 0.55, 1.25, 12.2, 0.7, [("Powerful oracles disagree, are uncalibrated, and conflate molecular effect with disease cause. POLARIS turns the two guardrails into mathematics.", 16, SLATE, False, "Calibri")])
for i, (tt, bd, col) in enumerate([
    ("Guardrail 1 — functional ≠ pathogenic",
     "Posterior factorizes into molecular-function × disease-linkage. Rarity is INVERTED for common-variant traits (replicated in 2 diseases).", TEAL),
    ("Guardrail 2 — distrust any single oracle",
     "Each channel is a noisy annotator; a latent-truth EM learns its reliability. Channels are largely independent (Marchenko–Pastur).", AMBER)]):
    card = s.shapes.add_shape(RR, I(0.55 + i * 6.25), I(2.35), I(5.9), I(3.4))
    card.fill.solid(); card.fill.fore_color.rgb = col; card.line.fill.background(); card.adjustments[0] = 0.06
    tf = card.text_frame; tf.word_wrap = True
    tf.margin_left = tf.margin_right = Pt(16); tf.margin_top = Pt(16)
    p = tf.paragraphs[0]; r = p.add_run(); r.text = tt; r.font.size = Pt(19); r.font.bold = True; r.font.color.rgb = WHITE; r.font.name = "Cambria"
    p2 = tf.add_paragraph(); p2.space_before = Pt(10); r2 = p2.add_run(); r2.text = bd; r2.font.size = Pt(15); r2.font.color.rgb = WHITE; r2.font.name = "Calibri"
transition(s)

# 3-9 figure slides
fig_slide("The pipeline — six steps, entirely from public data", "fig_pipeline",
          "Ranked variant → gene hypotheses with an auditable mechanism and a calibrated probability")
fig_slide("Every engine built from scratch and validated", "fig_methods_validation",
          "SuSiE fine-mapping, Bayesian coloc, reliability EM, and the information-theoretic motif logo — all validated",
          "from-scratch statistical & biophysical engines")
fig_slide("Calibrated triangulation in type 2 diabetes", "fig_results_main",
          "The EM learns conservation is most reliable; causal variants are conserved but NOT rare; channels independent; ECE = 0.005")
fig_slide("Generalization — the same behavior transfers to a second disease", "fig_generalization",
          "Conservation enrichment, rarity inversion, reliability ordering, and gene nomination all replicate in CAD",
          "T2D → coronary artery disease")
fig_slide("Vignette — FTO rs1421085 → ARID5B → IRX3", "fig_vignette_fto",
          "One variant, three gene calls. POLARIS recovers the textbook distal mechanism from first principles",
          "biophysics → dynamics → calibrated hypothesis")
fig_slide("The evidence landscape of top hypotheses", "fig_evidence_heatmap",
          "18 of 22 gene calls correct; the misses are the known-hard distal/imprinted loci, each with a transparent mechanism")
fig_slide("Statistical rigor — bootstrap, permutation, abstention", "fig_robustness",
          "phyloP dominance is significant (perm p = 0.013); confidence-stratified abstention raises gene precision to 0.89")

# 10 — key numbers (dark)
s = prs.slides.add_slide(BLANK); bg(s, NAVY)
textbox(s, 0.7, 0.5, 12, 0.8, [("By the numbers", 34, WHITE, True, "Cambria")])
stats = [(f"{m['gene_top1_POLARIS(=L2G)']*100:.0f}%", "gene top-1 (control loci)"),
         ("569 + 260", "fine-mapped variants (T2D+CAD)"),
         (f"{m['calib_ece']}", "calibration error (ECE)"),
         (f"0.68 / 0.87", "conservation AUC (T2D / CAD)"),
         (f"{m['oracle_spectrum']['n_channels']-m['oracle_spectrum']['n_signal_dims']}/{m['oracle_spectrum']['n_channels']}", "independent evidence channels"),
         ("p = 0.013", "phyloP reliability (permutation)")]
for i, (v, l) in enumerate(stats):
    cx = 0.7 + (i % 3) * 4.15; cy = 1.7 + (i // 3) * 2.4
    card = s.shapes.add_shape(RR, I(cx), I(cy), I(3.85), I(2.1))
    card.fill.solid(); card.fill.fore_color.rgb = INK; card.line.color.rgb = TEAL; card.line.width = Pt(1); card.adjustments[0] = 0.08
    tf = card.text_frame; tf.word_wrap = True; tf.vertical_anchor = MSO_ANCHOR.MIDDLE
    p = tf.paragraphs[0]; p.alignment = PP_ALIGN.CENTER; r = p.add_run(); r.text = v
    r.font.size = Pt(40); r.font.bold = True; r.font.color.rgb = TEAL if i % 2 == 0 else AMBER; r.font.name = "Cambria"
    p2 = tf.add_paragraph(); p2.alignment = PP_ALIGN.CENTER; r2 = p2.add_run(); r2.text = l
    r2.font.size = Pt(13); r2.font.color.rgb = MIST; r2.font.name = "Calibri"
transition(s)

# 11 — conclusion (dark)
s = prs.slides.add_slide(BLANK); bg(s, NAVY)
textbox(s, 1.0, 1.5, 11.3, 1.0, [("POLARIS", 48, WHITE, True, "Cambria")], PP_ALIGN.CENTER)
for i, t in enumerate([
    "Credibility = calibrated convergence across independent, fallible evidence — not any single oracle",
    "Transparent biophysics recovers textbook mechanisms (ARID5B, SORT1) from first principles",
    "Oracle-agnostic: black boxes like AlphaGenome can be calibrated, not blindly trusted"]):
    textbox(s, 1.6, 3.0 + i * 0.78, 10.1, 0.7, [("•  " + t, 18, MIST, False, "Calibri")], PP_ALIGN.LEFT)
textbox(s, 1.0, 6.3, 11.3, 0.6, [("Candidate A · GWAS hit → mechanistic hypothesis", 15, TEAL, False, "Calibri")], PP_ALIGN.CENTER)
transition(s)

out = config.ROOT / "presentation" / "POLARIS_deck.pptx"
prs.save(str(out))
print("saved", out, "| slides:", len(prs.slides._sldIdLst))
